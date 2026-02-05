#!/usr/bin/env python3
"""Create JSONL batch requests with support for multiple document formats.

⚠️  IMPORTANT: This script is for SIMPLE, UNIFORM batch processing only.

USE THIS WHEN:
  ✓ Same prompt for ALL files
  ✓ Same model for ALL files
  ✓ No conditional logic needed

DO NOT USE for complex cases (different prompts/models per file).
For advanced scenarios, generate custom code. See SKILL.md.

CONFIGURATION:
- Edit prompt.txt for your task instructions
- Set DOUBLEWORD_MODEL in .env.dw (default: Qwen3-VL-235B)
- Adjust MAX_TOKENS based on expected output length
- Use --extensions to filter file types

SUPPORTED FORMATS: PDF, DOCX, PPTX, ODP, TXT, MD, TSV, CSV, XLS, XLSX

WORKFLOW:
  1. python create_batch.py --input-dir /path/to/files/
  2. python submit_batch.py
  3. python poll_and_process.py
"""

import json
import os
import argparse
from pathlib import Path
import glob
import tomllib
import hashlib
from dotenv import load_dotenv
from datetime import datetime
import sys

# Load configuration from config.toml and .env.dw
def load_config():
    """Load config from dw_batch/config.toml and merge with .env.dw secrets."""
    # Load TOML config (non-secrets)
    config_path = Path(__file__).parent / 'config.toml'
    if not config_path.exists():
        print(f"Error: Configuration file not found: {config_path}")
        print("Please ensure config.toml exists")
        sys.exit(1)

    with open(config_path, 'rb') as f:
        config = tomllib.load(f)

    # Load .env.dw for secrets
    env_path = Path(__file__).parent / '.env.dw'
    load_dotenv(dotenv_path=env_path)

    # Check for required secret
    auth_token = os.getenv('DOUBLEWORD_AUTH_TOKEN')
    if not auth_token:
        print("="*60)
        print("ERROR: DOUBLEWORD_AUTH_TOKEN not found")
        print("="*60)
        print("Please ensure you have:")
        print("1. Created .env.dw file from .env.dw.sample")
        print("2. Added your DOUBLEWORD_AUTH_TOKEN to .env.dw")
        print("="*60)
        sys.exit(1)

    return config, auth_token

config, auth_token = load_config()

# Parse command line arguments
parser = argparse.ArgumentParser(
    description='Create JSONL batch requests from documents',
    formatter_class=argparse.RawDescriptionHelpFormatter,
    epilog='''
Examples:
  # Process all files in default directory (data/papers/)
  python create_batch.py

  # Process specific files
  python create_batch.py --files paper1.pdf paper2.csv data.xlsx

  # Process all files in a custom directory
  python create_batch.py --input-dir /path/to/documents/

  # Process only specific file types
  python create_batch.py --input-dir /path/to/data/ --extensions csv xlsx
'''
)
parser.add_argument(
    '--files',
    nargs='+',
    metavar='FILE',
    help='Specific file paths to process'
)
parser.add_argument(
    '--input-dir',
    metavar='DIR',
    help='Directory to scan for documents (default: data/papers/)'
)
parser.add_argument(
    '--extensions',
    nargs='+',
    metavar='EXT',
    help='Only process these extensions (e.g., csv tsv xlsx). Default: all supported'
)
parser.add_argument(
    '--output-dir',
    metavar='DIR',
    required=True,
    help='Output directory for results (REQUIRED - agent must pass absolute path to project root)'
)
parser.add_argument(
    '--logs-dir',
    metavar='DIR',
    help='Directory for logs and batch files (default: {output-dir}/logs)'
)
parser.add_argument(
    '--dry-run',
    action='store_true',
    help='Estimate costs without creating batch file (recommended for large jobs)'
)
parser.add_argument(
    '--skip-existing',
    action='store_true',
    help='Skip files that already have output summaries (opt-in, checks filename + prompt hash)'
)
parser.add_argument(
    '--force',
    action='store_true',
    help='Force batch creation even if cost thresholds are exceeded (use with caution)'
)

args = parser.parse_args()

# Read prompt template
with open('prompt.txt', 'r') as f:
    prompt_template = f.read()

# Substitute word count from config
word_count = str(config['output']['summary_word_count'])
prompt_template = prompt_template.replace('{WORD_COUNT}', word_count)

# Compute prompt hash for skip-existing check
prompt_hash = hashlib.md5(prompt_template.encode()).hexdigest()[:8]

# Get config values
model = config['models']['default_model']
max_tokens = config['output']['max_tokens']
chat_endpoint = config['api']['chat_completions_endpoint']

# Print configuration being used
print("="*60)
print("BATCH REQUEST CONFIGURATION")
print("="*60)
print(f"Prompt template: prompt.txt")
print(f"Model: {model}")
print(f"Max tokens: {max_tokens}")
print(f"Word count target: {word_count}")
print(f"Auth token: ...{auth_token[-4:]}")
print("="*60)
print()

# Collect files based on arguments
all_supported_extensions = ['*.pdf', '*.txt', '*.md', '*.docx', '*.pptx', '*.odp', '*.tsv', '*.csv', '*.xls', '*.xlsx']

if args.extensions:
    # Filter to only requested extensions
    supported_extensions = [f'*.{ext.lstrip("*.")}' for ext in args.extensions]
else:
    supported_extensions = all_supported_extensions

all_files = []

if args.files:
    # Use specific files provided
    all_files = [str(Path(f).resolve()) for f in args.files]
    print(f"Processing {len(all_files)} specified file(s)\n")
elif args.input_dir:
    # Scan custom directory
    input_dir = Path(args.input_dir)
    if not input_dir.exists():
        print(f"Error: Directory '{args.input_dir}' does not exist")
        exit(1)
    for ext in supported_extensions:
        all_files.extend(glob.glob(str(input_dir / ext)))
    all_files.sort()
    print(f"Found {len(all_files)} files in {args.input_dir}\n")
else:
    # Default: scan data/papers directory
    for ext in supported_extensions:
        all_files.extend(glob.glob(f'data/papers/{ext}'))
    all_files.sort()
    print(f"Found {len(all_files)} files in data/papers/\n")

if not all_files:
    print("No files found to process. Exiting.")
    exit(0)

# Skip-existing setup
if args.skip_existing:
    print("\n" + "="*60)
    print("⚠️  SKIP-EXISTING MODE ENABLED")
    print("="*60)
    print("This will skip files that already have output summaries.")
    print("Files are matched by: filename + prompt hash")
    print(f"Current prompt hash: {prompt_hash}")
    print()
    print("⚠️  WARNING: Use this carefully!")
    print("  - If you changed the prompt, old outputs may be stale")
    print("  - Prompt hash changes = all files reprocessed")
    print("  - This is opt-in to avoid accidental re-runs")
    print("="*60)
    print()

    # Determine output directory to check for existing files
    output_dir = Path(args.output_dir)

    # Find existing summary files
    existing_summaries = {}
    if output_dir.exists():
        for summary_file in output_dir.glob('*_summary_*.md'):
            # Extract base filename (everything before _summary_)
            base_name = summary_file.stem.split('_summary_')[0]
            existing_summaries[base_name] = summary_file

        if existing_summaries:
            print(f"Found {len(existing_summaries)} existing summary file(s) in {output_dir}")
        else:
            print(f"No existing summaries found in {output_dir}")
    else:
        print(f"Output directory {output_dir} does not exist yet")
    print()

# Detect what file types we actually have
detected_extensions = set(Path(f).suffix.lower() for f in all_files)
print(f"Detected file types: {', '.join(sorted(detected_extensions))}\n")

# CONDITIONAL IMPORTS - Only import what we need based on detected file types
print("Loading required libraries...")
imports_loaded = []

if any(ext in detected_extensions for ext in ['.pdf']):
    from pypdf import PdfReader
    import pdfplumber
    imports_loaded.append('PDF libraries')

if '.docx' in detected_extensions:
    from docx import Document
    imports_loaded.append('DOCX')

if '.pptx' in detected_extensions:
    from pptx import Presentation
    imports_loaded.append('PPTX')

if '.odp' in detected_extensions:
    from odf.opendocument import load as load_odf
    from odf.text import P
    from odf.draw import Frame
    imports_loaded.append('ODP')

if any(ext in detected_extensions for ext in ['.csv', '.tsv', '.xls', '.xlsx']):
    import pandas as pd
    imports_loaded.append('Pandas (Excel/CSV/TSV)')

print(f"✓ Loaded: {', '.join(imports_loaded)}\n")

# Initialize tracking
requests = []
failed_files = []
extraction_stats = {}
total_input_chars = 0  # Track for dry-run cost estimation

# ============================================================================
# EXTRACTION FUNCTIONS
# ============================================================================

def extract_text_pypdf(pdf_path):
    """Try pypdf first (faster)."""
    with open(pdf_path, 'rb') as f:
        reader = PdfReader(f)
        text = '\n'.join(page.extract_text() for page in reader.pages)
        return text, len(reader.pages)

def extract_text_pdfplumber(pdf_path):
    """Fallback to pdfplumber (more robust but slower)."""
    with pdfplumber.open(pdf_path) as pdf:
        text = '\n'.join((page.extract_text() or '') for page in pdf.pages)
        return text, len(pdf.pages)

def extract_from_text(file_path):
    """Extract text from .txt or .md files."""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        text = f.read()
        return text, 1

def extract_from_docx(file_path):
    """Extract text from .docx files."""
    doc = Document(file_path)
    paragraphs = [para.text for para in doc.paragraphs]
    text = '\n'.join(paragraphs)
    # Estimate pages (rough: 500 words per page)
    word_count_local = len(text.split())
    pages = max(1, word_count_local // 500)
    return text, pages

def extract_from_pptx(file_path):
    """Extract text from .pptx files."""
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    text = '\n'.join(text_runs)
    return text, len(prs.slides)

def extract_from_odp(file_path):
    """Extract text from .odp files."""
    doc = load_odf(file_path)
    text_runs = []
    # Extract all text paragraphs
    for paragraph in doc.getElementsByType(P):
        text_content = ''.join(node.data for node in paragraph.childNodes if hasattr(node, 'data'))
        if text_content.strip():
            text_runs.append(text_content)
    text = '\n'.join(text_runs)
    # Count frames as slide estimate
    frames = doc.getElementsByType(Frame)
    pages = max(1, len(frames))
    return text, pages

def extract_from_excel(file_path):
    """Extract data from Excel files (.xls, .xlsx) as formatted text."""
    df = pd.read_excel(file_path, sheet_name=None)  # Read all sheets

    text_parts = []
    total_rows = 0

    for sheet_name, sheet_df in df.items():
        # Clean the dataframe (drop completely empty rows/columns)
        sheet_df = sheet_df.dropna(how='all').dropna(axis=1, how='all')

        if not sheet_df.empty:
            text_parts.append(f"### Sheet: {sheet_name}\n")
            # Convert to tab-separated format for better readability
            text_parts.append(sheet_df.to_csv(sep='\t', index=False))
            text_parts.append("\n")
            total_rows += len(sheet_df)

    text = '\n'.join(text_parts)
    pages = max(1, total_rows // 50)  # Estimate: 50 rows per "page"
    return text, pages

def extract_from_csv_tsv(file_path):
    """Extract data from CSV/TSV files as formatted text."""
    # Auto-detect delimiter
    try:
        df = pd.read_csv(file_path, sep=None, engine='python', encoding='utf-8', on_bad_lines='skip')
    except Exception:
        # Fallback: try common delimiters explicitly
        for sep in [',', '\t', ';', '|']:
            try:
                df = pd.read_csv(file_path, sep=sep, encoding='utf-8', on_bad_lines='skip')
                break
            except Exception:
                continue
        else:
            raise ValueError("Could not parse CSV/TSV file with common delimiters")

    # Clean the dataframe
    df = df.dropna(how='all').dropna(axis=1, how='all')

    # Convert to tab-separated format
    text = df.to_csv(sep='\t', index=False)
    pages = max(1, len(df) // 50)  # Estimate: 50 rows per "page"
    return text, pages

# ============================================================================
# PROCESS FILES
# ============================================================================

for idx, file_path in enumerate(all_files, 1):
    print(f"[{idx}/{len(all_files)}] Processing {file_path}...")

    # Check if we should skip this file (if --skip-existing enabled)
    if args.skip_existing:
        file_stem = Path(file_path).stem
        # Sanitize filename same way as custom_id generation
        safe_filename = file_stem.replace('%', '_').replace(' ', '_').replace('&', 'and')[:55]

        if safe_filename in existing_summaries:
            print(f"  ⏭️  Skipped (existing summary: {existing_summaries[safe_filename].name})")
            continue

    text = None
    pages = 0
    extraction_method = None
    file_extension = Path(file_path).suffix.lower()

    try:
        # Route to appropriate extraction method based on file type
        if file_extension == '.pdf':
            # Try pypdf first (faster), fallback to pdfplumber
            try:
                text, pages = extract_text_pypdf(file_path)
                extraction_method = 'pypdf'
            except (KeyError, Exception) as e:
                if 'bbox' in str(e) or isinstance(e, KeyError):
                    print(f"  ⚠ pypdf failed ({e}), trying pdfplumber...")
                    text, pages = extract_text_pdfplumber(file_path)
                    extraction_method = 'pdfplumber'
                else:
                    raise

        elif file_extension == '.docx':
            text, pages = extract_from_docx(file_path)
            extraction_method = 'docx'

        elif file_extension == '.pptx':
            text, pages = extract_from_pptx(file_path)
            extraction_method = 'pptx'

        elif file_extension == '.odp':
            text, pages = extract_from_odp(file_path)
            extraction_method = 'odp'

        elif file_extension in ['.txt', '.md']:
            text, pages = extract_from_text(file_path)
            extraction_method = 'txt'

        elif file_extension in ['.xls', '.xlsx']:
            text, pages = extract_from_excel(file_path)
            extraction_method = 'excel'

        elif file_extension in ['.csv', '.tsv']:
            text, pages = extract_from_csv_tsv(file_path)
            extraction_method = 'csv/tsv'

        else:
            print(f"  ⚠ Unsupported file type: {file_extension}")
            failed_files.append((file_path, f"unsupported file type: {file_extension}"))
            continue

        # Track extraction stats
        extraction_stats[extraction_method] = extraction_stats.get(extraction_method, 0) + 1

    except Exception as e:
        print(f"\n{'!'*60}")
        print(f"ERROR DURING PRE-PROCESSING: {Path(file_path).name}")
        print(f"{'!'*60}")
        print(f"File: {file_path}")
        print(f"Extension: {file_extension}")
        print(f"Error type: {type(e).__name__}")
        print(f"Error message: {str(e)}")
        print(f"{'!'*60}\n")
        failed_files.append((file_path, f"{type(e).__name__}: {str(e)}"))
        continue

    # Skip if no meaningful text extracted
    if not text or len(text.strip()) < 100:
        print(f"  ⚠ Skipped (insufficient text: {len(text)} chars)")
        failed_files.append((file_path, "insufficient text"))
        continue

    print(f"  ✓ Extracted {len(text)} characters from {pages} pages [{extraction_method}]")

    # Track input chars for dry-run cost estimation (prompt + document text)
    total_input_chars += len(text) + len(prompt_template)

    # Create batch request with sanitized custom_id
    # Remove special chars from filename for custom_id (max 64 chars including 'summary-' prefix)
    safe_filename = Path(file_path).stem.replace('%', '_').replace(' ', '_').replace('&', 'and')[:55]

    request = {
        "custom_id": f"summary-{safe_filename}",
        "method": "POST",
        "url": chat_endpoint,
        "body": {
            "model": model,
            "messages": [
                {
                    "role": "user",
                    "content": f"{prompt_template}\n\nDocument text:\n{text}"
                }
            ],
            "max_tokens": max_tokens
        }
    }
    requests.append(request)

# ============================================================================
# DRY RUN MODE - ESTIMATE COSTS BEFORE CREATING BATCH
# ============================================================================

if args.dry_run:
    # Rough token estimates (4 chars ≈ 1 token for English text)
    estimated_input_tokens = total_input_chars // 4
    estimated_output_tokens = len(requests) * max_tokens

    # Cost estimation based on Doubleword pricing (Feb 2026)
    # From screenshot: Qwen3-VL-235B: $0.125/1M input (1h), Qwen3-VL-30B: $0.07/1M input (1h)
    if '235B' in model:
        cost_per_1m = 0.125  # 1h pricing
        model_display = "Qwen3-VL-235B (complex)"
    elif '30B' in model:
        cost_per_1m = 0.07  # 1h pricing
        model_display = "Qwen3-VL-30B (simple)"
    else:
        cost_per_1m = 0.10  # fallback estimate
        model_display = model

    # Simple estimate: same rate for input/output (conservative)
    estimated_cost = ((estimated_input_tokens + estimated_output_tokens) / 1_000_000) * cost_per_1m

    completion_window = config['batch']['completion_window']

    # Check cost thresholds (unless --force flag is used)
    max_input_tokens = config['safety']['max_input_tokens']
    max_output_tokens = config['safety']['max_output_tokens']

    threshold_exceeded = (
        estimated_input_tokens > max_input_tokens or
        estimated_output_tokens > max_output_tokens
    )

    print("\n" + "="*60)
    print("DRY RUN - COST ESTIMATION")
    print("="*60)
    print(f"Files to process: {len(requests)}")
    print(f"Files failed/skipped: {len(failed_files)}")
    print(f"Model: {model_display}")
    print(f"Completion window: {completion_window}")
    print()
    print(f"Estimated input tokens: ~{estimated_input_tokens:,}")
    print(f"Max output tokens: {estimated_output_tokens:,} ({len(requests)} × {max_tokens})")
    print(f"Total tokens: ~{estimated_input_tokens + estimated_output_tokens:,}")
    print()
    print(f"Estimated cost: ${estimated_cost:.4f}")
    print()
    print("⚠️  This is a ROUGH ESTIMATE:")
    print("  - Token count is approximate (4 chars ≈ 1 token)")
    print("  - Output cost assumes MAX_TOKENS per file (worst case)")
    print("  - Actual costs may be lower if responses are shorter")
    print("="*60)

    # Check thresholds and warn if exceeded
    if threshold_exceeded and not args.force:
        print("\n" + "!"*60)
        print("⚠️  COST THRESHOLD EXCEEDED - REVIEW REQUIRED")
        print("!"*60)
        print("Safety thresholds from config.toml:")
        print(f"  Max input tokens: {max_input_tokens:,}")
        print(f"  Max output tokens: {max_output_tokens:,}")
        print()
        print("Your batch exceeds these limits:")
        if estimated_input_tokens > max_input_tokens:
            print(f"  ✗ Input tokens: {estimated_input_tokens:,} (exceeds {max_input_tokens:,})")
        else:
            print(f"  ✓ Input tokens: {estimated_input_tokens:,} (within limit)")
        if estimated_output_tokens > max_output_tokens:
            print(f"  ✗ Output tokens: {estimated_output_tokens:,} (exceeds {max_output_tokens:,})")
        else:
            print(f"  ✓ Output tokens: {estimated_output_tokens:,} (within limit)")
        print()
        print("⚠️  RECOMMENDED ACTIONS:")
        print("  1. Reduce file count (use --files to select subset)")
        print("  2. Lower MAX_TOKENS in config.toml")
        print("  3. Use smaller model (Qwen3-VL-30B instead of 235B)")
        print("  4. Split into multiple smaller batches")
        print()
        print("To proceed anyway, add --force flag (use with caution)")
        print("="*60)
        sys.exit(1)
    elif threshold_exceeded and args.force:
        print("\n" + "⚠️ "*30)
        print("WARNING: Cost thresholds exceeded but proceeding due to --force flag")
        print(f"  Input tokens: {estimated_input_tokens:,} (limit: {max_input_tokens:,})")
        print(f"  Output tokens: {estimated_output_tokens:,} (limit: {max_output_tokens:,})")
        print("⚠️ "*30)

    print("\nTo proceed with batch creation, remove --dry-run flag")
    print("="*60)
    sys.exit(0)

# ============================================================================
# CHECK COST THRESHOLDS (NON-DRY-RUN MODE)
# ============================================================================

# Calculate token estimates for threshold check
estimated_input_tokens = total_input_chars // 4
estimated_output_tokens = len(requests) * max_tokens

# Check safety thresholds
max_input_tokens = config['safety']['max_input_tokens']
max_output_tokens = config['safety']['max_output_tokens']

threshold_exceeded = (
    estimated_input_tokens > max_input_tokens or
    estimated_output_tokens > max_output_tokens
)

if threshold_exceeded and not args.force:
    print("\n" + "!"*60)
    print("⚠️  COST THRESHOLD EXCEEDED - ABORTING")
    print("!"*60)
    print("Safety thresholds from config.toml:")
    print(f"  Max input tokens: {max_input_tokens:,}")
    print(f"  Max output tokens: {max_output_tokens:,}")
    print()
    print("Your batch exceeds these limits:")
    if estimated_input_tokens > max_input_tokens:
        print(f"  ✗ Input tokens: {estimated_input_tokens:,} (exceeds {max_input_tokens:,})")
    else:
        print(f"  ✓ Input tokens: {estimated_input_tokens:,} (within limit)")
    if estimated_output_tokens > max_output_tokens:
        print(f"  ✗ Output tokens: {estimated_output_tokens:,} (exceeds {max_output_tokens:,})")
    else:
        print(f"  ✓ Output tokens: {estimated_output_tokens:,} (within limit)")
    print()
    print("⚠️  RECOMMENDED ACTIONS:")
    print("  1. Run with --dry-run first to estimate costs")
    print("  2. Reduce file count (use --files to select subset)")
    print("  3. Lower MAX_TOKENS in config.toml")
    print("  4. Use smaller model (Qwen3-VL-30B instead of 235B)")
    print("  5. Split into multiple smaller batches")
    print()
    print("To proceed anyway, add --force flag (use with caution)")
    print("="*60)
    sys.exit(1)
elif threshold_exceeded and args.force:
    print("\n" + "⚠️ "*30)
    print("WARNING: Cost thresholds exceeded but proceeding due to --force flag")
    print(f"  Input tokens: {estimated_input_tokens:,} (limit: {max_input_tokens:,})")
    print(f"  Output tokens: {estimated_output_tokens:,} (limit: {max_output_tokens:,})")
    print("⚠️ "*30)
    print()

# ============================================================================
# SAVE BATCH REQUEST FILE TO LOGS FOLDER
# ============================================================================

# Determine output and logs directories
output_dir = Path(args.output_dir)
logs_dir = Path(args.logs_dir) if args.logs_dir else (output_dir / 'logs')
logs_dir.mkdir(parents=True, exist_ok=True)

# Print absolute paths for clarity
print(f"\nOutput directory: {output_dir.resolve()}")
print(f"Logs directory: {logs_dir.resolve()}")

# Write JSONL file with timestamp to logs folder
timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
output_file = logs_dir / f'batch_requests_{timestamp}.jsonl'
with open(output_file, 'w') as f:
    for req in requests:
        f.write(json.dumps(req) + '\n')

print(f"\n{'='*60}")
print(f"✓ Created {output_file} with {len(requests)} requests")
print(f"\nExtraction methods used:")
for method, count in sorted(extraction_stats.items()):
    print(f"  {method}: {count} files")

# Log failed files to error log if any
if failed_files:
    error_log_file = logs_dir / f'batch_errors_{timestamp}.log'
    with open(error_log_file, 'w') as f:
        f.write(f"Batch Creation Error Log\n")
        f.write(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
        f.write(f"Total files failed: {len(failed_files)}\n")
        f.write(f"{'='*60}\n\n")
        for path, reason in failed_files:
            f.write(f"File: {path}\n")
            f.write(f"Reason: {reason}\n")
            f.write(f"{'-'*60}\n")

    print(f"\n⚠ Failed to process {len(failed_files)} files:")
    for path, reason in failed_files:
        print(f"  - {Path(path).name}: {reason}")
    print(f"\n✓ Error details saved to: {error_log_file}")

print(f"\n{'='*60}")
print("NEXT STEPS:")
print(f"1. Review the batch file: {output_file}")
print(f"2. Submit the batch: python submit_batch.py {output_file}")
print("3. Monitor progress: python poll_and_process.py")
print("="*60)
