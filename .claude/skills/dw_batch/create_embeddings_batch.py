#!/usr/bin/env python3
"""Create JSONL batch requests for embeddings with support for multiple document formats.

⚠️  IMPORTANT: This script is for SIMPLE, UNIFORM embedding generation only.

USE THIS WHEN:
  ✓ Same embedding model for ALL files
  ✓ Need to embed full document text
  ✓ No conditional logic needed

CONFIGURATION:
- Set DOUBLEWORD_EMBEDDING_MODEL in .env.dw
- Use --extensions to filter file types
- Use --chunk-size to split long documents (default: no chunking)

SUPPORTED FORMATS: PDF, DOCX, PPTX, ODP, TXT, MD, TSV, CSV, XLS, XLSX

WORKFLOW:
  1. python create_embeddings_batch.py --input-dir /path/to/files/
  2. python submit_batch.py
  3. python poll_and_process.py
"""

import json
import os
import argparse
from pathlib import Path
import glob
from dotenv import load_dotenv
from datetime import datetime

# Parse command line arguments
parser = argparse.ArgumentParser(
    description='Create JSONL batch requests for embeddings from documents',
    formatter_class=argparse.RawDescriptionHelpFormatter,
    epilog='''
Examples:
  # Process all files in default directory (data/)
  python create_embeddings_batch.py

  # Process specific files
  python create_embeddings_batch.py --files paper1.pdf paper2.csv data.xlsx

  # Process all files in a custom directory
  python create_embeddings_batch.py --input-dir /path/to/documents/

  # Process only specific file types
  python create_embeddings_batch.py --input-dir /path/to/data/ --extensions csv xlsx

  # Chunk long documents (useful for context limits)
  python create_embeddings_batch.py --chunk-size 2000
'''
)
parser.add_argument(
    '--files',
    nargs='+',
    metavar='FILE',
    help='Specific file paths to process'
)
parser.add_argument(
    '--input-dir',
    metavar='DIR',
    help='Directory to scan for documents (default: data/)'
)
parser.add_argument(
    '--extensions',
    nargs='+',
    metavar='EXT',
    help='Only process these extensions (e.g., csv tsv xlsx). Default: all supported'
)
parser.add_argument(
    '--chunk-size',
    type=int,
    metavar='TOKENS',
    help='Split documents into chunks of this many tokens (default: no chunking)'
)

args = parser.parse_args()

# Load environment variables
load_dotenv()

# Get embedding model from environment (default to a common model)
embedding_model = os.getenv('DOUBLEWORD_EMBEDDING_MODEL', 'BAAI/bge-en-icl')

# Print configuration
print("="*60)
print("EMBEDDINGS BATCH REQUEST CONFIGURATION")
print("="*60)
print(f"Model: {embedding_model}")
if args.chunk_size:
    print(f"Chunk size: {args.chunk_size} tokens (~{args.chunk_size * 0.75:.0f} words)")
else:
    print(f"Chunk size: No chunking (full document)")
print("="*60)
print()

# Collect files based on arguments
all_supported_extensions = ['*.pdf', '*.txt', '*.md', '*.docx', '*.pptx', '*.odp', '*.tsv', '*.csv', '*.xls', '*.xlsx']

if args.extensions:
    # Filter to only requested extensions
    extensions = [f'*.{ext.lstrip("*.")}' for ext in args.extensions]
else:
    extensions = all_supported_extensions

all_files = []

if args.files:
    # Use explicitly specified files
    all_files = args.files
    print(f"Processing {len(all_files)} explicitly specified files")
else:
    # Scan directory for matching files
    input_dir = args.input_dir if args.input_dir else '../../data/'
    input_path = Path(input_dir)

    if not input_path.exists():
        print(f"ERROR: Input directory does not exist: {input_path}")
        exit(1)

    print(f"Scanning {input_path} for files...")
    for ext in extensions:
        all_files.extend(glob.glob(str(input_path / '**' / ext), recursive=True))

    print(f"Found {len(all_files)} files matching extensions: {', '.join(extensions)}")

if not all_files:
    print("ERROR: No files found to process!")
    exit(1)

# ============================================================================
# CONDITIONAL IMPORTS - Only import what we need based on detected file types
# ============================================================================

detected_extensions = set(Path(f).suffix.lower() for f in all_files)
print(f"\nDetected file types: {', '.join(sorted(detected_extensions))}")
print("Loading required libraries...")

if any(ext in detected_extensions for ext in ['.pdf']):
    from pypdf import PdfReader
    import pdfplumber

if any(ext in detected_extensions for ext in ['.docx']):
    from docx import Document

if any(ext in detected_extensions for ext in ['.pptx']):
    from pptx import Presentation

if any(ext in detected_extensions for ext in ['.odp']):
    from odf.opendocument import load as load_odf
    from odf.text import P
    from odf.draw import Frame

if any(ext in detected_extensions for ext in ['.csv', '.tsv', '.xls', '.xlsx']):
    import pandas as pd

print("✓ Libraries loaded\n")

# ============================================================================
# TEXT EXTRACTION FUNCTIONS (reused from create_batch.py)
# ============================================================================

requests = []
failed_files = []
extraction_stats = {}

def extract_text_pypdf(file_path):
    """Extract text from PDF using pypdf (fast)."""
    reader = PdfReader(file_path)
    text_parts = []
    for page in reader.pages:
        text_parts.append(page.extract_text())
    text = '\n'.join(text_parts)
    return text, len(reader.pages)

def extract_text_pdfplumber(file_path):
    """Extract text from PDF using pdfplumber (better formatting)."""
    text_parts = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text_parts.append(page.extract_text() or '')
    text = '\n'.join(text_parts)
    return text, len(text_parts)

def extract_from_docx(file_path):
    """Extract text from .docx files."""
    doc = Document(file_path)
    text_parts = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text_parts.append(paragraph.text)
    text = '\n'.join(text_parts)
    word_count_local = len(text.split())
    pages = max(1, word_count_local // 500)
    return text, pages

def extract_from_text(file_path):
    """Extract text from .txt or .md files."""
    with open(file_path, 'r', encoding='utf-8') as f:
        text = f.read()
    word_count_local = len(text.split())
    pages = max(1, word_count_local // 500)
    return text, pages

def extract_from_pptx(file_path):
    """Extract text from .pptx files."""
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    text = '\n'.join(text_runs)
    return text, len(prs.slides)

def extract_from_odp(file_path):
    """Extract text from .odp files."""
    doc = load_odf(file_path)
    text_runs = []
    for paragraph in doc.getElementsByType(P):
        text_content = ''.join(node.data for node in paragraph.childNodes if hasattr(node, 'data'))
        if text_content.strip():
            text_runs.append(text_content)
    text = '\n'.join(text_runs)
    frames = doc.getElementsByType(Frame)
    pages = max(1, len(frames))
    return text, pages

def extract_from_excel(file_path):
    """Extract data from Excel files (.xls, .xlsx) as formatted text."""
    try:
        df = pd.read_excel(file_path, sheet_name=None)
    except Exception as e:
        raise Exception(f"Failed to read Excel file: {str(e)}")

    text_parts = []
    total_rows = 0

    for sheet_name, sheet_df in df.items():
        sheet_df = sheet_df.dropna(how='all').dropna(axis=1, how='all')

        if not sheet_df.empty:
            text_parts.append(f"### Sheet: {sheet_name}\n")
            text_parts.append(sheet_df.to_csv(sep='\t', index=False))
            text_parts.append("\n")
            total_rows += len(sheet_df)

    text = '\n'.join(text_parts)
    pages = max(1, total_rows // 50)
    return text, pages

def extract_from_csv_tsv(file_path):
    """Extract data from CSV/TSV files as formatted text."""
    try:
        df = pd.read_csv(file_path, sep=None, engine='python', encoding='utf-8', on_bad_lines='skip')
    except Exception as e:
        # Fallback: try common delimiters explicitly
        for sep in [',', '\t', ';', '|']:
            try:
                df = pd.read_csv(file_path, sep=sep, encoding='utf-8', on_bad_lines='skip')
                break
            except Exception:
                continue
        else:
            raise Exception(f"Could not parse CSV/TSV file with common delimiters: {str(e)}")

    df = df.dropna(how='all').dropna(axis=1, how='all')
    text = df.to_csv(sep='\t', index=False)
    pages = max(1, len(df) // 50)
    return text, pages

def chunk_text(text, chunk_size_tokens):
    """Split text into chunks of approximately chunk_size_tokens.

    Uses simple word-based splitting: 1 token ≈ 0.75 words
    """
    words = text.split()
    words_per_chunk = int(chunk_size_tokens * 0.75)

    chunks = []
    for i in range(0, len(words), words_per_chunk):
        chunk = ' '.join(words[i:i + words_per_chunk])
        chunks.append(chunk)

    return chunks

# ============================================================================
# PROCESS FILES
# ============================================================================

for idx, file_path in enumerate(all_files, 1):
    print(f"[{idx}/{len(all_files)}] Processing {file_path}...")

    text = None
    pages = 0
    extraction_method = None
    file_extension = Path(file_path).suffix.lower()

    try:
        # Route to appropriate extraction method based on file type
        if file_extension == '.pdf':
            try:
                text, pages = extract_text_pypdf(file_path)
                extraction_method = 'pypdf'
            except (KeyError, Exception) as e:
                if 'bbox' in str(e) or isinstance(e, KeyError):
                    print(f"  ⚠ pypdf failed ({e}), trying pdfplumber...")
                    text, pages = extract_text_pdfplumber(file_path)
                    extraction_method = 'pdfplumber'
                else:
                    raise

        elif file_extension == '.docx':
            text, pages = extract_from_docx(file_path)
            extraction_method = 'docx'

        elif file_extension == '.pptx':
            text, pages = extract_from_pptx(file_path)
            extraction_method = 'pptx'

        elif file_extension == '.odp':
            text, pages = extract_from_odp(file_path)
            extraction_method = 'odp'

        elif file_extension in ['.txt', '.md']:
            text, pages = extract_from_text(file_path)
            extraction_method = 'txt'

        elif file_extension in ['.xls', '.xlsx']:
            text, pages = extract_from_excel(file_path)
            extraction_method = 'excel'

        elif file_extension in ['.csv', '.tsv']:
            text, pages = extract_from_csv_tsv(file_path)
            extraction_method = 'csv/tsv'

        else:
            print(f"  ⚠ Unsupported file type: {file_extension}")
            failed_files.append((file_path, f"unsupported file type: {file_extension}"))
            continue

        extraction_stats[extraction_method] = extraction_stats.get(extraction_method, 0) + 1

    except Exception as e:
        print(f"\n{'!'*60}")
        print(f"ERROR DURING PRE-PROCESSING: {Path(file_path).name}")
        print(f"{'!'*60}")
        print(f"File: {file_path}")
        print(f"Extension: {file_extension}")
        print(f"Error type: {type(e).__name__}")
        print(f"Error message: {str(e)}")
        print(f"{'!'*60}\n")
        failed_files.append((file_path, f"{type(e).__name__}: {str(e)}"))
        continue

    # Skip if no meaningful text extracted
    if not text or len(text.strip()) < 10:
        print(f"  ⚠ Skipped (insufficient text: {len(text)} chars)")
        failed_files.append((file_path, "insufficient text"))
        continue

    print(f"  ✓ Extracted {len(text)} characters from {pages} pages [{extraction_method}]")

    # Handle chunking if requested
    if args.chunk_size:
        chunks = chunk_text(text, args.chunk_size)
        print(f"  → Split into {len(chunks)} chunks")
    else:
        chunks = [text]

    # Create batch requests for each chunk
    safe_filename = Path(file_path).stem.replace('%', '_').replace(' ', '_').replace('&', 'and')[:55]

    for chunk_idx, chunk in enumerate(chunks):
        # Create custom_id that includes chunk number if chunked
        if len(chunks) > 1:
            custom_id = f"embed-{safe_filename}-chunk{chunk_idx+1}"
        else:
            custom_id = f"embed-{safe_filename}"

        request = {
            "custom_id": custom_id,
            "method": "POST",
            "url": "/v1/embeddings",
            "body": {
                "model": embedding_model,
                "input": chunk
            }
        }
        requests.append(request)

# ============================================================================
# SAVE BATCH REQUEST FILE TO LOGS FOLDER
# ============================================================================

logs_dir = Path('../../dw_batch_output/logs')
logs_dir.mkdir(parents=True, exist_ok=True)

timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
output_file = logs_dir / f'batch_requests_embeddings_{timestamp}.jsonl'
with open(output_file, 'w') as f:
    for req in requests:
        f.write(json.dumps(req) + '\n')

print(f"\n{'='*60}")
print(f"✓ Created {output_file} with {len(requests)} embedding requests")
print(f"\nExtraction methods used:")
for method, count in sorted(extraction_stats.items()):
    print(f"  {method}: {count} files")

if failed_files:
    print(f"\n⚠ Failed to process {len(failed_files)} files:")
    for path, reason in failed_files:
        print(f"  - {Path(path).name}: {reason}")

print(f"\n{'='*60}")
print("NEXT STEPS:")
print(f"1. Review the batch file: {output_file}")
print(f"2. Submit the batch: python submit_batch.py")
print("3. Monitor progress: python poll_and_process.py")
print(f"{'='*60}")
